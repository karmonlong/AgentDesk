"""
PPT 生成工具
基于 Gemini API 生成 PPT 大纲和幻灯片图片
"""
import os
import json
import base64
import re
import requests
from typing import Dict, List, Optional, Any
from datetime import datetime


def generate_presentation_outline(
    topic: str,
    document_content: Optional[str] = None,
    complexity_level: str = "专业",
    visual_style: str = "现代简约",
    language: str = "简体中文",
    slide_count: int = 5
) -> Dict[str, Any]:
    """
    生成 PPT 大纲
    
    Args:
        topic: 主题或内容
        document_content: 文档内容（可选）
        complexity_level: 复杂度级别（通用/专业/学术/行政高管）
        visual_style: 视觉风格（现代简约/商务科技/创意艺术/深色模式/自然清新）
        language: 语言
        slide_count: 幻灯片数量
    
    Returns:
        包含 outline 和 sources 的字典
    """
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        return {
            "success": False,
            "error": "未配置 GEMINI_API_KEY"
        }
    
    base_url = os.getenv("GEMINI_BASE_URL", "https://generativelanguage.googleapis.com/v1beta")
    model = os.getenv("GEMINI_MODEL", "gemini-3-pro-preview")
    api_url = f"{base_url.rstrip('/')}/models/{model}:generateContent"
    
    # 风格指令
    style_instructions = {
        "现代简约": "Style: Minimalist, plenty of white space, sans-serif typography, soft pastel accents. Professional and airy.",
        "商务科技": "Style: Fortune 500 business style. Navy blues, greys, structured layouts, official look, subtle grid lines.",
        "创意艺术": "Style: Bold colors, artistic shapes, dynamic composition. High energy.",
        "深色模式": "Style: Dark background (slate/black), bright neon or white text, sleek and modern.",
        "自然清新": "Style: Organic shapes, green and earth tones, soft lighting."
    }
    style_instr = style_instructions.get(visual_style, style_instructions["现代简约"])
    
    # 构建提示词
    is_article = len(topic) > 250 or document_content is not None
    
    base_instruction = f"""
你是一个专业的演示文稿设计师。
任务：创建一个包含 {slide_count} 张幻灯片的演示文稿大纲。

受众级别：{complexity_level}
语言：{language}（重要：JSON 内容、标题和正文必须使用 {language}）

输出规则：
1. 创建恰好 {slide_count} 张幻灯片
2. 第 1 张必须是标题页
3. 最后一张必须是总结/结论页
4. 严格以 JSON 数组格式返回输出

JSON 格式：
[
  {{
    "title": "幻灯片标题（使用 {language}）",
    "content": "幻灯片要点内容（最多30个词，使用 {language}）",
    "visualDescription": "为 AI 图像生成器创建此幻灯片背景和布局的详细提示词。包括文本占位符的位置。{style_instr}"
  }}
]
"""
    
    if is_article:
        prompt_text = f"""
{base_instruction}

分析提供的内容（文本和/或附加文档）来创建演示文稿。
"""
        if topic:
            prompt_text += f'\n\n要分析的内容：\n"{topic[:20000]}"'
        if document_content:
            prompt_text += f'\n\n附加文档内容：\n"{document_content[:20000]}"'
    else:
        prompt_text = f"""
{base_instruction}

任务：研究主题 "{topic}" 并创建演示文稿。
**使用 Google Search 查找准确的事实。**
"""
    
    # 构建请求
    parts = [{"text": prompt_text}]
    
    # 如果有文档内容，可以添加到 parts（如果是 PDF 可以添加 inlineData）
    # 这里简化处理，只使用文本
    
    payload = {
        "contents": [{"parts": parts}]
    }
    
    # 如果不是文章模式，可以添加 Google Search 工具
    if not is_article:
        payload["tools"] = [{"googleSearch": {}}]
    
    try:
        print(f"[PPTGen] 生成演示文稿大纲...")
        print(f"  主题: {topic[:50]}...")
        print(f"  幻灯片数量: {slide_count}")
        print(f"  复杂度: {complexity_level}")
        print(f"  风格: {visual_style}")
        
        resp = requests.post(
            api_url,
            headers={
                "Content-Type": "application/json",
                "x-goog-api-key": api_key,
            },
            json=payload,
            timeout=120,
        )
        
        if resp.status_code != 200:
            return {
                "success": False,
                "error": f"API 返回错误: HTTP {resp.status_code}",
                "details": resp.text[:500]
            }
        
        result = resp.json()
        text = ""
        
        # 提取文本内容
        candidates = result.get("candidates", [])
        if candidates:
            content = candidates[0].get("content", {})
            parts_list = content.get("parts", [])
            for part in parts_list:
                if "text" in part:
                    text += part["text"]
        
        # 提取来源
        sources = []
        grounding_metadata = candidates[0].get("groundingMetadata", {}) if candidates else {}
        chunks = grounding_metadata.get("groundingChunks", [])
        for chunk in chunks:
            web = chunk.get("web", {})
            if web.get("uri") and web.get("title"):
                sources.append({
                    "title": web["title"],
                    "url": web["uri"]
                })
        
        # 去重来源
        unique_sources = list({s["url"]: s for s in sources}.values())
        
        # 解析 JSON 大纲
        outline = []
        try:
            import re
            json_match = re.search(r'\[[\s\S]*\]', text)
            if json_match:
                outline = json.loads(json_match.group(0))
            else:
                # 尝试直接解析整个文本
                outline = json.loads(text)
        except Exception as e:
            print(f"[PPTGen] JSON 解析失败: {e}")
            # 返回错误
            return {
                "success": False,
                "error": f"无法解析大纲 JSON: {str(e)}",
                "raw_text": text[:500]
            }
        
        print(f"[PPTGen] ✅ 成功生成 {len(outline)} 张幻灯片大纲")
        
        return {
            "success": True,
            "outline": outline,
            "sources": unique_sources
        }
        
    except Exception as e:
        print(f"[PPTGen] ❌ 生成失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }


def generate_slide_image(
    slide_outline: Dict[str, str],
    visual_style: str = "现代简约"
) -> Dict[str, Any]:
    """
    生成单张幻灯片图片
    
    Args:
        slide_outline: 包含 title, content, visualDescription 的字典
        visual_style: 视觉风格
    
    Returns:
        包含图片 base64 数据的字典
    """
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        return {
            "success": False,
            "error": "未配置 GEMINI_API_KEY"
        }
    
    base_url = os.getenv("GEMINI_BASE_URL", "https://generativelanguage.googleapis.com/v1beta")
    model = "gemini-3-pro-image-preview"  # 使用图像生成模型
    api_url = f"{base_url.rstrip('/')}/models/{model}:generateContent"
    
    # 构建图像生成提示词
    prompt = f"""
创建一张高质量的 16:9 演示文稿幻灯片图片。

幻灯片上的文本（清晰渲染此文本）：
标题：{slide_outline.get('title', '')}
正文：{slide_outline.get('content', '')}

设计说明：
{slide_outline.get('visualDescription', '')}

重要：
- 图片必须包含上面提供的文本
- 文本必须清晰、大号、专业
- 使用 16:9 宽高比布局
- 不要包含混乱或乱码文本
"""
    
    payload = {
        "contents": [{
            "parts": [{"text": prompt}]
        }],
        "generationConfig": {
            "responseModalities": ["IMAGE"]
        }
    }
    
    try:
        print(f"[PPTGen] 生成幻灯片图片: {slide_outline.get('title', '')[:30]}...")
        
        resp = requests.post(
            api_url,
            headers={
                "Content-Type": "application/json",
                "x-goog-api-key": api_key,
            },
            json=payload,
            timeout=120,
        )
        
        if resp.status_code != 200:
            return {
                "success": False,
                "error": f"API 返回错误: HTTP {resp.status_code}",
                "details": resp.text[:500]
            }
        
        result = resp.json()
        
        # 提取图片数据
        candidates = result.get("candidates", [])
        if candidates:
            content = candidates[0].get("content", {})
            parts_list = content.get("parts", [])
            for part in parts_list:
                inline_data = part.get("inlineData") or part.get("inline_data")
                if inline_data and inline_data.get("data"):
                    print(f"[PPTGen] ✅ 成功生成幻灯片图片")
                    return {
                        "success": True,
                        "image_base64": inline_data["data"],
                        "mime_type": inline_data.get("mimeType", "image/png")
                    }
        
        return {
            "success": False,
            "error": "未找到图片数据",
            "response": result
        }
        
    except Exception as e:
        print(f"[PPTGen] ❌ 生成图片失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }


def create_pdf_from_slides(
    slides: List[Dict[str, Any]],
    output_path: str,
    title: str = "演示文稿"
) -> Dict[str, Any]:
    """
    从幻灯片数据创建 PDF 文件
    
    Args:
        slides: 幻灯片列表，每个包含 title, content, image_base64
        output_path: 输出文件路径
        title: 演示文稿标题
    
    Returns:
        成功/失败状态
    """
    try:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        from io import BytesIO
        import base64
        from PIL import Image
        
        # 使用横向 A4 (16:9 比例)
        page_width, page_height = landscape(A4)
        
        pdf = canvas.Canvas(output_path, pagesize=(page_width, page_height))
        
        for idx, slide_data in enumerate(slides):
            if slide_data.get("image_base64"):
                try:
                    # 解码 base64 图片
                    img_data = base64.b64decode(slide_data["image_base64"])
                    img = Image.open(BytesIO(img_data))
                    
                    # 调整图片大小以适应页面（保持 16:9 比例）
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    # 计算适合页面的尺寸
                    if aspect_ratio > page_width / page_height:
                        # 图片更宽，以宽度为准
                        draw_width = page_width
                        draw_height = page_width / aspect_ratio
                    else:
                        # 图片更高，以高度为准
                        draw_height = page_height
                        draw_width = page_height * aspect_ratio
                    
                    # 居中显示
                    x = (page_width - draw_width) / 2
                    y = (page_height - draw_height) / 2
                    
                    # 将 PIL Image 转换为 ReportLab 可用的格式
                    img_buffer = BytesIO()
                    img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    
                    # 绘制图片
                    pdf.drawImage(ImageReader(img_buffer), x, y, width=draw_width, height=draw_height)
                    
                except Exception as img_err:
                    print(f"[PPTGen] 添加图片失败: {img_err}")
                    # 如果图片加载失败，至少显示文本
                    pdf.setFont("Helvetica-Bold", 24)
                    pdf.drawString(50, page_height - 100, slide_data.get("title", f"幻灯片 {idx + 1}"))
                    pdf.setFont("Helvetica", 14)
                    pdf.drawString(50, page_height - 150, slide_data.get("content", ""))
            else:
                # 没有图片，只显示文本
                pdf.setFont("Helvetica-Bold", 24)
                pdf.drawString(50, page_height - 100, slide_data.get("title", f"幻灯片 {idx + 1}"))
                pdf.setFont("Helvetica", 14)
                pdf.drawString(50, page_height - 150, slide_data.get("content", ""))
            
            # 添加新页面（除了最后一张）
            if idx < len(slides) - 1:
                pdf.showPage()
        
        pdf.save()
        print(f"[PPTGen] ✅ PDF 文件已保存: {output_path}")
        
        return {
            "success": True,
            "output_path": output_path
        }
        
    except ImportError as e:
        # 检查缺少哪个库
        missing_lib = str(e).split("'")[1] if "'" in str(e) else "未知库"
        return {
            "success": False,
            "error": f"需要安装 {missing_lib}: pip install {missing_lib}",
            "hint": "生成 PDF 需要: pip install reportlab pillow"
        }
    except Exception as e:
        print(f"[PPTGen] ❌ 创建 PDF 失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }


def create_pptx_from_slides(
    slides: List[Dict[str, Any]],
    output_path: str,
    title: str = "演示文稿"
) -> Dict[str, Any]:
    """
    从幻灯片数据创建 PPTX 文件
    
    Args:
        slides: 幻灯片列表，每个包含 title, content, image_base64
        output_path: 输出文件路径
        title: 演示文稿标题
    
    Returns:
        成功/失败状态
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        import base64
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for idx, slide_data in enumerate(slides):
            # 选择布局
            if idx == 0:
                # 第一张使用标题布局
                slide_layout = prs.slide_layouts[0]  # Title slide
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加标题
            if slide_layout == prs.slide_layouts[0]:
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", "")
            else:
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", "")
                
                # 添加内容
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.text = slide_data.get("content", "")
            
            # 如果有图片，添加图片
            if slide_data.get("image_base64"):
                try:
                    img_data = base64.b64decode(slide_data["image_base64"])
                    img_stream = BytesIO(img_data)
                    
                    # 添加图片到幻灯片
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(10)
                    height = Inches(7.5)
                    slide.shapes.add_picture(img_stream, left, top, width, height)
                except Exception as img_err:
                    print(f"[PPTGen] 添加图片失败: {img_err}")
        
        prs.save(output_path)
        print(f"[PPTGen] ✅ PPTX 文件已保存: {output_path}")
        
        return {
            "success": True,
            "output_path": output_path
        }
        
    except ImportError:
        # 如果没有 python-pptx，返回提示
        return {
            "success": False,
            "error": "需要安装 python-pptx: pip install python-pptx",
            "hint": "可以使用其他方式生成 PPT，或返回图片列表"
        }
    except Exception as e:
        print(f"[PPTGen] ❌ 创建 PPTX 失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }

